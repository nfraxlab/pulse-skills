#!/usr/bin/env python3
"""
PowerPoint Generator Script
Generates clean, professional presentations using python-pptx
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def create_presentation(data):
    """
    Create a PowerPoint presentation from structured data.
    
    Expected data format:
    {
        "title": "Presentation Title",
        "subtitle": "Optional subtitle",
        "author": "Author name",
        "slides": [
            {
                "type": "title",  # or "content" or "section"
                "title": "Slide Title",
                "content": ["Bullet 1", "Bullet 2", ...] or "Single text block"
            }
        ],
        "output_file": "output.pptx"
    }
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.get("title", "Untitled Presentation")
    subtitle_text = data.get("subtitle", "")
    if data.get("author"):
        subtitle_text += f"\n{data['author']}" if subtitle_text else data["author"]
    subtitle.text = subtitle_text
    
    # Style title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
    
    # Create content slides
    for slide_data in data.get("slides", []):
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "section":
            # Section header slide
            section_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(section_layout)
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
            title.text_frame.paragraphs[0].font.size = Pt(40)
            title.text_frame.paragraphs[0].font.bold = True
            
        else:
            # Standard content slide
            bullet_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_layout)
            title = slide.shapes.title
            title.text = slide_data.get("title", "")
            
            # Style title
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 41, 55)
            
            # Add content
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                content = slide_data.get("content", [])
                if isinstance(content, str):
                    content = [content]
                
                for idx, item in enumerate(content):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(75, 85, 99)
                    p.space_after = Pt(12)
    
    # Save presentation
    output_file = data.get("output_file", "presentation.pptx")
    prs.save(output_file)
    print(f"✓ Presentation created: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file


def main():
    """Main entry point - reads JSON from stdin or file"""
    if len(sys.argv) > 1:
        # Read from file
        input_file = sys.argv[1]
        with open(input_file, 'r') as f:
            data = json.load(f)
    else:
        # Read from stdin
        data = json.load(sys.stdin)
    
    create_presentation(data)


if __name__ == "__main__":
    main()
